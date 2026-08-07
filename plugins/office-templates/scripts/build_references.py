#!/usr/bin/env python3
"""Generate minimal reference Office templates for 灵犀 office-templates plugin."""

from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1] / "templates"


def make_letterhead(path: Path) -> None:
    d = Document()
    t = d.add_paragraph("LINGXI")
    t.runs[0].bold = True
    t.runs[0].font.size = Pt(16)
    t.runs[0].font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
    sub = d.add_paragraph("商务信笺模板 · Minimal Letterhead")
    sub.runs[0].font.size = Pt(10)
    sub.runs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    d.add_paragraph("")
    d.add_paragraph("收件人：")
    d.add_paragraph("主题：")
    d.add_paragraph("")
    d.add_paragraph("在此填写正文。保持简洁、礼貌、可执行。")
    d.add_paragraph("")
    d.add_paragraph("此致")
    d.add_paragraph("敬礼")
    d.add_paragraph("")
    d.add_paragraph("署名 / 日期")
    d.save(path)


def make_memo(path: Path, title: str, sections) -> None:
    d = Document()
    d.add_heading(title, level=1)
    for name, hint in sections:
        d.add_heading(name, level=2)
        d.add_paragraph(hint)
    d.save(path)


def make_deck(path: Path, title: str, slides) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title
    try:
        s.placeholders[1].text = "灵犀办公模板"
    except Exception:
        pass
    body_layout = prs.slide_layouts[1]
    for st, bullets in slides:
        sl = prs.slides.add_slide(body_layout)
        sl.shapes.title.text = st
        tf = sl.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
    prs.save(path)


thin = Border(
    left=Side(style="thin", color="D1D5DB"),
    right=Side(style="thin", color="D1D5DB"),
    top=Side(style="thin", color="D1D5DB"),
    bottom=Side(style="thin", color="D1D5DB"),
)
header_fill = PatternFill("solid", fgColor="1F2937")
header_font = Font(color="FFFFFF", bold=True)


def style_header(ws, row, cols):
    for c in range(1, cols + 1):
        cell = ws.cell(row, c)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")
        cell.border = thin


def make_dashboard(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Dashboard"
    ws["A1"] = "分析看板"
    ws["A1"].font = Font(bold=True, size=16)
    headers = ["指标", "本周", "上周", "环比", "备注"]
    for i, h in enumerate(headers, 1):
        ws.cell(3, i, h)
    style_header(ws, 3, 5)
    data = [
        ("活跃用户", 1200, 1100, "=B4/C4-1", ""),
        ("转化率", 0.082, 0.075, "=B5/C5-1", ""),
        ("收入", 86000, 79000, "=B6/C6-1", ""),
    ]
    for r, row in enumerate(data, 4):
        for c, v in enumerate(row, 1):
            ws.cell(r, c, v).border = thin
    for cell in ("D4", "D5", "D6"):
        ws[cell].number_format = "0.0%"
    chart = BarChart()
    chart.title = "本周 vs 上周"
    chart.add_data(Reference(ws, min_col=2, min_row=3, max_col=3, max_row=6), titles_from_data=True)
    chart.set_categories(Reference(ws, min_col=1, min_row=4, max_row=6))
    ws.add_chart(chart, "A9")
    wb.save(path)


def make_budget(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws["A1"] = "财务预算"
    ws["A1"].font = Font(bold=True, size=16)
    headers = ["科目", "Q1", "Q2", "Q3", "Q4", "合计"]
    for i, h in enumerate(headers, 1):
        ws.cell(3, i, h)
    style_header(ws, 3, 6)
    for r, name in enumerate(["人力", "市场", "研发", "运营", "其他"], 4):
        ws.cell(r, 1, name).border = thin
        for c in range(2, 6):
            ws.cell(r, c, 0).border = thin
        ws.cell(r, 6, f"=SUM(B{r}:E{r})").border = thin
    ws.cell(9, 1, "总计").font = Font(bold=True)
    for c, col in enumerate("BCDEF", 2):
        ws.cell(9, c, f"=SUM({col}4:{col}8)").border = thin
    wb.save(path)


def make_tracker(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Tracker"
    ws["A1"] = "项目跟踪"
    ws["A1"].font = Font(bold=True, size=16)
    headers = ["任务", "Owner", "状态", "优先级", "开始", "截止", "进度%", "备注"]
    for i, h in enumerate(headers, 1):
        ws.cell(3, i, h)
    style_header(ws, 3, 8)
    sample = ["示例任务", "未分配", "进行中", "P1", "", "", 0.3, ""]
    for c, v in enumerate(sample, 1):
        cell = ws.cell(4, c, v)
        cell.border = thin
        if c == 7:
            cell.number_format = "0%"
    for col, w in zip("ABCDEFGH", [22, 12, 10, 10, 12, 12, 10, 20]):
        ws.column_dimensions[col].width = w
    wb.save(path)


def make_three_statement(path: Path) -> None:
    wb = Workbook()
    a = wb.active
    a.title = "Assumptions"
    a["A1"] = "三表预测 · 假设"
    a["A1"].font = Font(bold=True, size=14)
    a["A3"] = "增长率"
    a["B3"] = 0.05
    a["B3"].number_format = "0.0%"
    a["A4"] = "税率"
    a["B4"] = 0.25
    a["B4"].number_format = "0.0%"

    is_ = wb.create_sheet("IncomeStatement")
    is_["A1"] = "损益表"
    is_["A1"].font = Font(bold=True, size=14)
    for i, h in enumerate(["项目", "Y1", "Y2", "Y3"], 1):
        is_.cell(3, i, h)
    style_header(is_, 3, 4)
    items = [
        ("收入", 1000, "=B4*(1+Assumptions!B3)", "=C4*(1+Assumptions!B3)"),
        ("成本", 400, "=B5*(1+Assumptions!B3)", "=C5*(1+Assumptions!B3)"),
        ("毛利", "=B4-B5", "=C4-C5", "=D4-D5"),
        ("费用", 200, 210, 220),
        ("税前利润", "=B6-B7", "=C6-C7", "=D6-D7"),
        ("所得税", "=B8*Assumptions!B4", "=C8*Assumptions!B4", "=D8*Assumptions!B4"),
        ("净利润", "=B8-B9", "=C8-C9", "=D8-D9"),
    ]
    for r, row in enumerate(items, 4):
        for c, v in enumerate(row, 1):
            is_.cell(r, c, v).border = thin

    bs = wb.create_sheet("BalanceSheet")
    bs["A1"] = "资产负债表（简）"
    bs["A3"] = "资产合计"
    bs["B3"] = 1500
    bs["A4"] = "负债合计"
    bs["B4"] = 600
    bs["A5"] = "所有者权益"
    bs["B5"] = "=B3-B4"

    cf = wb.create_sheet("CashFlow")
    cf["A1"] = "现金流量表（简）"
    cf["A3"] = "经营现金流"
    cf["B3"] = "=IncomeStatement!B10"
    cf["A4"] = "投资现金流"
    cf["B4"] = -100
    cf["A5"] = "筹资现金流"
    cf["B5"] = 50
    cf["A6"] = "净现金流"
    cf["B6"] = "=B3+B4+B5"
    wb.save(path)


def main() -> None:
    ROOT.mkdir(parents=True, exist_ok=True)
    make_letterhead(ROOT / "minimal-letterhead.docx")
    make_memo(
        ROOT / "strategy-memorandum.docx",
        "战略备忘录",
        [
            ("背景", "说明决策背景与约束。"),
            ("目标", "写清要达成的结果。"),
            ("选项", "列出可选方案与取舍。"),
            ("建议", "给出明确推荐。"),
            ("风险与里程碑", "关键风险与下一步。"),
        ],
    )
    make_memo(
        ROOT / "system-design.docx",
        "系统设计说明",
        [
            ("需求摘要", "功能与非功能需求。"),
            ("架构总览", "组件与边界。"),
            ("数据流", "关键路径与存储。"),
            ("接口", "主要 API / 事件。"),
            ("权衡与运维", "取舍、监控与发布。"),
        ],
    )
    make_memo(
        ROOT / "design-report.docx",
        "设计报告",
        [
            ("问题定义", "用户与痛点。"),
            ("设计目标", "成功标准。"),
            ("方案概述", "核心交互与视觉方向。"),
            ("验收要点", "可测标准。"),
        ],
    )
    make_deck(
        ROOT / "team-alignment.pptx",
        "团队对齐",
        [
            ("议程", ["背景", "目标", "优先级", "决策", "行动项"]),
            ("目标", ["本季度核心结果", "非目标"]),
            ("优先级", ["P0", "P1", "P2"]),
            ("决策与行动", ["已决事项", "Owner / 截止日期"]),
        ],
    )
    make_deck(
        ROOT / "project-kickoff.pptx",
        "项目启动",
        [
            ("为什么做", ["业务动机", "成功定义"]),
            ("范围", ["In scope", "Out of scope"]),
            ("里程碑", ["M1", "M2", "M3"]),
            ("角色与协作", ["RACI 要点", "沟通节奏"]),
        ],
    )
    make_deck(
        ROOT / "business-review.pptx",
        "业务复盘",
        [
            ("总览", ["关键指标", "结论一句话"]),
            ("进展", ["已完成", "进行中"]),
            ("问题与风险", ["阻塞", "缓解"]),
            ("下阶段计划", ["重点动作", "资源需求"]),
        ],
    )
    make_deck(
        ROOT / "simple-light.pptx",
        "简约演示",
        [
            ("要点一", ["说明", "证据"]),
            ("要点二", ["说明", "证据"]),
            ("总结", ["行动项"]),
        ],
    )
    make_dashboard(ROOT / "analytics-dashboard.xlsx")
    make_budget(ROOT / "financial-budget.xlsx")
    make_tracker(ROOT / "project-tracker.xlsx")
    make_three_statement(ROOT / "three-statement-forecast.xlsx")
    print("ok", sorted(p.name for p in ROOT.iterdir()))


if __name__ == "__main__":
    main()
